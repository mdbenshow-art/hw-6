# -*- coding: utf-8 -*-
"""
generate_presentation.py

Generates a beautiful, high-quality PowerPoint presentation (.pptx) based on the
CRISP-DM lifecycle for Multiple Linear Regression using the 50 Startups dataset.
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

# 1. Initialize Presentation with 16:9 widescreen
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Colors
c_bg = RGBColor(15, 23, 42)          # Slate-900
c_card_bg = RGBColor(30, 41, 59)     # Slate-800
c_border = RGBColor(51, 65, 85)       # Slate-700
c_cyan = RGBColor(6, 182, 212)       # Cyan-500 (Title Accent)
c_purple = RGBColor(168, 85, 247)     # Purple-500 (Highlight Accent)
c_white = RGBColor(241, 245, 249)     # Slate-100 (Primary Text)
c_muted = RGBColor(148, 163, 184)     # Slate-400 (Secondary Text)
c_red = RGBColor(239, 68, 68)         # Red-500 (Alert/Highlight)

def set_slide_background(slide, color):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_slide_header(slide, title_text):
    txBox = slide.shapes.add_textbox(Inches(0.7), Inches(0.4), Inches(11.933), Inches(0.8))
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.bold = True
    p.font.size = Pt(28)
    p.font.color.rgb = c_cyan
    p.font.name = "Microsoft JhengHei"

def add_card(slide, left, top, width, height, title, content_lines, bg_color=c_card_bg, border_color=c_border, title_color=c_cyan, text_color=c_white):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg_color
    shape.line.color.rgb = border_color
    shape.line.width = Pt(1.5)
    
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.25)
    tf.margin_right = Inches(0.25)
    tf.margin_top = Inches(0.25)
    tf.margin_bottom = Inches(0.25)
    
    p0 = tf.paragraphs[0]
    p0.text = title
    p0.font.bold = True
    p0.font.size = Pt(16)
    p0.font.color.rgb = title_color
    p0.font.name = "Microsoft JhengHei"
    p0.space_after = Pt(10)
    p0.alignment = PP_ALIGN.LEFT
    
    for line in content_lines:
        p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(12)
        p.font.color.rgb = text_color
        p.font.name = "Microsoft JhengHei"
        p.space_after = Pt(5)
        p.alignment = PP_ALIGN.LEFT

# ==========================================
# SLIDE 1: Cover Slide
# ==========================================
slide = prs.slides.add_slide(prs.slide_layouts[6]) # blank layout
set_slide_background(slide, c_bg)

# Main Title & Subtitle in a single textframe
title_box = slide.shapes.add_textbox(Inches(1.0), Inches(2.2), Inches(11.333), Inches(3.0))
tf = title_box.text_frame
tf.word_wrap = True

p_main = tf.paragraphs[0]
p_main.text = "基於 CRISP-DM 生命週期的\n多元線性迴歸教學簡報"
p_main.font.bold = True
p_main.font.size = Pt(44)
p_main.font.color.rgb = c_cyan
p_main.font.name = "Microsoft JhengHei"
p_main.space_after = Pt(20)
p_main.alignment = PP_ALIGN.CENTER

p_sub = tf.add_paragraph()
p_sub.text = "—— 以 50 Startups 新創利潤預測與特徵選取最佳化 (Sweet Spot) 為核心案例"
p_sub.font.size = Pt(18)
p_sub.font.color.rgb = c_purple
p_sub.font.name = "Microsoft JhengHei"
p_sub.space_after = Pt(40)
p_sub.alignment = PP_ALIGN.CENTER

p_meta = tf.add_paragraph()
p_meta.text = "資料科學與機器學習教學小組  •  2026年6月"
p_meta.font.size = Pt(12)
p_meta.font.color.rgb = c_muted
p_meta.font.name = "Microsoft JhengHei"
p_meta.alignment = PP_ALIGN.CENTER

# ==========================================
# SLIDE 2: CRISP-DM Overview
# ==========================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_background(slide, c_bg)
add_slide_header(slide, "CRISP-DM 數據挖掘流程生命週期")

# Draw 6 circular/rect process cards horizontally
card_width = Inches(1.8)
card_height = Inches(4.5)
y_pos = Inches(1.8)
start_x = Inches(0.7)
gap = Inches(0.2)

phases = [
    ("1. 商業理解\n(Business)", ["目標：定義業務需求與利潤預測指標。", "核心：為創投 (VC) 篩選高獲利公司並降低投資風險。"]),
    ("2. 數據理解\n(Data)", ["目標：收集與探索 50 Startups 資料集。", "核心：分析支出特徵與利潤 (Profit) 的相關性。"]),
    ("3. 數據準備\n(Preparation)", ["目標：數據清洗與工程化。", "核心：執行數值標準化與地區 One-Hot 編碼。"]),
    ("4. 模型建置\n(Modeling)", ["目標：建立與調優多元線性迴歸。", "核心：對比線性模型、Ridge 懲罰與隨機森林。"]),
    ("5. 模型評估\n(Evaluation)", ["目標：評估模型表現與尋找最佳點。", "核心：利用 10 組演算法找出 k=2 最佳特徵點。"]),
    ("6. 模型部署\n(Deployment)", ["目標：封裝 Pipeline 並投入生產環境。", "核心：導出為 pkl，並整合 FastAPI 作即時預測。"])
]

for idx, (title, content) in enumerate(phases):
    x_pos = start_x + idx * (card_width + gap)
    # Highlight modeling & evaluation
    border_color = c_purple if "評估" in title or "建置" in title else c_border
    title_color = c_purple if "評估" in title or "建置" in title else c_cyan
    add_card(slide, x_pos, y_pos, card_width, card_height, title, content, border_color=border_color, title_color=title_color)

# ==========================================
# SLIDE 3: Business Understanding
# ==========================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_background(slide, c_bg)
add_slide_header(slide, "階段一：商業理解 (Business Understanding)")

# Description text box
desc_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.4), Inches(11.933), Inches(1.0))
tf = desc_box.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "本階段旨在確立新創公司的評估架構。對於創投機構而言，盲目燒錢的企業往往伴隨極高風險。我們需要找出哪些支出項目能切實為企業帶來利潤增幅，從而量化早期投資的風險。"
p.font.size = Pt(14)
p.font.color.rgb = c_white
p.font.name = "Microsoft JhengHei"
p.space_after = Pt(10)

# Add 3 key cards
add_card(slide, Inches(0.7), Inches(2.6), Inches(3.7), Inches(3.8), "商業痛點與問題", [
    "• 早期投資缺乏定量指標：",
    "  難以量化研發、行政與行銷的投入產出比率。",
    "• 行業資訊不透明：",
    "  如何有效辨識被投企業的組織效率與冗餘成本？",
    "• 地理位置效益不明："
    "  各州份 (加州、紐約州、佛州) 經營紅利是否存在顯著差異？"
])

add_card(slide, Inches(4.8), Inches(2.6), Inches(3.7), Inches(3.8), "機器學習分析目標", [
    "• 目標變數 (Target)：",
    "  預測新創公司的總利潤 (Profit)。",
    "• 解釋變數 (Features)：",
    "  研發支出、行政管理支出、市場行銷支出與公司所屬州份。",
    "• 關鍵指標：",
    "  解釋特徵邊際貢獻（即迴歸係數的大小與正負向）。"
])

add_card(slide, Inches(8.9), Inches(2.6), Inches(3.7), Inches(3.8), "模型驗證成功準則", [
    "• 統計顯著性指標 (R²)：",
    "  測試集模型決定係數 $R^2 \\ge 90\\%$，代表能準確捕捉九成以上的利潤變異。",
    "• 預測絕對誤差 (MAE) 限制：",
    "  平均預測絕對誤差 $MAE \\le \\$7,000$ 美元，確保決策精度在可靠區間。"
])

# ==========================================
# SLIDE 4: Data Understanding
# ==========================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_background(slide, c_bg)
add_slide_header(slide, "階段二：數據理解 (Data Understanding)")

# Text box
desc_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.3), Inches(11.933), Inches(1.0))
tf = desc_box.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "本階段探索 50 家美國早期新創公司資料集。我們重點關注各支出特徵的分佈、有無缺失值，以及它們與總利潤的關聯性。"
p.font.size = Pt(14)
p.font.color.rgb = c_white
p.font.name = "Microsoft JhengHei"

# Add a table for feature analysis
rows, cols = 5, 4
left, top_t, width_t, height_t = Inches(0.7), Inches(2.2), Inches(6.8), Inches(4.0)
table_shape = slide.shapes.add_table(rows, cols, left, top_t, width_t, height_t)
table = table_shape.table

# Set Column Widths
table.columns[0].width = Inches(1.8)
table.columns[1].width = Inches(2.2)
table.columns[2].width = Inches(1.3)
table.columns[3].width = Inches(1.5)

headers = ["欄位名稱", "業務描述", "資料型態", "與利潤相關性"]
for i, head in enumerate(headers):
    cell = table.cell(0, i)
    cell.text = head
    cell.fill.solid()
    cell.fill.fore_color.rgb = RGBColor(17, 24, 39)
    p = cell.text_frame.paragraphs[0]
    p.font.bold = True
    p.font.size = Pt(13)
    p.font.color.rgb = c_cyan
    p.font.name = "Microsoft JhengHei"
    p.alignment = PP_ALIGN.CENTER

table_data = [
    ["R&D Spend", "研發領域支出(核心技術)", "數值型 (USD)", "0.9729 (極高)"],
    ["Marketing Spend", "市場推廣與宣傳支出", "數值型 (USD)", "0.7478 (高度)"],
    ["Administration", "行政支出(辦公租金等)", "數值型 (USD)", "0.2007 (微弱)"],
    ["State", "公司註冊所在地(紐約/加州/佛州)", "類別型 (String)", "地理效應分析"]
]

for row_idx, row_val in enumerate(table_data):
    for col_idx, val in enumerate(row_val):
        cell = table.cell(row_idx + 1, col_idx)
        cell.text = val
        cell.fill.solid()
        cell.fill.fore_color.rgb = c_card_bg
        p = cell.text_frame.paragraphs[0]
        p.font.size = Pt(11)
        p.font.color.rgb = c_white
        p.font.name = "Microsoft JhengHei"
        p.alignment = PP_ALIGN.CENTER

# Add 2 summary cards on the right
add_card(slide, Inches(8.0), Inches(2.2), Inches(4.6), Inches(1.9), "關鍵數據分佈特徵", [
    "• 樣本總量：50 筆，無缺失值。",
    "• 平均利潤：$112,012 美元，標準差約 $40,306。",
    "• 最顯著特徵：研發支出 (R&D Spend) 與利潤的相關性接近 0.97，是最核心的利潤驅動引擎。"
])

add_card(slide, Inches(8.0), Inches(4.3), Inches(4.6), Inches(1.9), "地理州別分析說明", [
    "• 地理分佈：紐約、加州、佛羅里達州三地分佈均勻。",
    "• 佛羅里達州的新創公司在 Profit 上平均值微幅領先，需在後續的特徵工程中進一步轉化以量化地理影響。"
])

# ==========================================
# SLIDE 5: Data Preparation
# ==========================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_background(slide, c_bg)
add_slide_header(slide, "階段三：數據準備 (Data Preparation)")

desc_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.3), Inches(11.933), Inches(1.0))
tf = desc_box.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "特徵準備是模型表現的基石。為了解決數據量綱（特徵量級）不一以及類別型特徵的處理，我們在此實施標準的前處理工程。"
p.font.size = Pt(14)
p.font.color.rgb = c_white
p.font.name = "Microsoft JhengHei"

# 3 Columns for Data Prep
add_card(slide, Inches(0.7), Inches(2.3), Inches(3.7), Inches(4.2), "① 金額特徵標準化\n(StandardScaler)", [
    "• 問題背景：",
    "  研發、行政與行銷支出在金額量級上有顯著差異。",
    "• 解決方案：",
    "  使用 scikit-learn 的 `StandardScaler` 將數值型金額欄位轉換為「均值為 0，方差為 1」的標準常態分佈特徵。",
    "• 優點：",
    "  消除量綱影響，提升迴歸模型收斂速度，使係數更具解釋性。"
])

add_card(slide, Inches(4.8), Inches(2.3), Inches(3.7), Inches(4.2), "② 類別型特徵編碼\n(One-Hot Encoding)", [
    "• 問題背景：",
    "  State (州別) 是字串類別，無法直接傳入數學方程。",
    "• 解決方案：",
    "  使用 `OneHotEncoder(drop='first')` 進行獨熱編碼。",
    "• 避免多重共線性：",
    "  以加州 (California) 為基準對照組 (Dropped)，產生 `State_Florida` 與 `State_New York` 兩個虛擬特徵。"
])

add_card(slide, Inches(8.9), Inches(2.3), Inches(3.7), Inches(4.2), "③ 前處理 Pipeline 封裝\n(ColumnTransformer)", [
    "• 整體 Pipeline 組裝：",
    "  將金額標準化與州別獨熱編碼封裝進單一 `ColumnTransformer` 架構中。",
    "• 避免數據洩漏 (Data Leakage)：",
    "  前處理只在訓練集上 `fit_transform`，測試集僅進行 `transform`，保證數據驗證的科學性。"
])

# ==========================================
# SLIDE 6: Modeling
# ==========================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_background(slide, c_bg)
add_slide_header(slide, "階段四：模型建置 (Modeling)")

desc_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.3), Inches(11.933), Inches(1.0))
tf = desc_box.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "本階段實施多模型並行策略。我們建立了基礎線性模型、帶有防過擬合機制的正則化模型以及具備非線性擬合能力的樹集成模型進行交叉對比。"
p.font.size = Pt(14)
p.font.color.rgb = c_white
p.font.name = "Microsoft JhengHei"

# 3 Modeling Cards
add_card(slide, Inches(0.7), Inches(2.3), Inches(3.7), Inches(4.2), "① 多元線性迴歸\n(Linear Regression)", [
    "• 核心特色：",
    "  直接估計各項支出與地理虛擬特徵的線性權重。",
    "• 業務價值：",
    "  具備極高的解釋力，可以直接提供「投入 1 美元能帶來多少邊際利潤」的明確指標。",
    "• 前提假設：",
    "  假設特徵與目標值之間存在高度線性關係。"
])

add_card(slide, Inches(4.8), Inches(2.3), Inches(3.7), Inches(4.2), "② 脊迴歸正則化\n(Ridge Regression)", [
    "• 核心特色：",
    "  在線性損失函數中加入 L2 正則化懲罰項 (L2 Penalty)。",
    "• 業務價值：",
    "  當特徵之間存在一定程度的共線性（如研發與行銷支出通常伴隨增長）時，Ridge 能夠對係數進行收縮，穩定預測結果，防止模型過擬合。"
])

add_card(slide, Inches(8.9), Inches(2.3), Inches(3.7), Inches(4.2), "③ 隨機森林與參數優化\n(Random Forest & GridSearch)", [
    "• 核心特色：",
    "  非線性多決策樹集成算法。",
    "• 業務價值：",
    "  不受多重共線性的干擾，能自動捕捉非線性特徵交互作用。",
    "• 超參數調優：",
    "  使用 5 折交叉驗證 (`GridSearchCV`) 調整樹數量與深度，獲取泛化性能最優的估計器。"
])

# ==========================================
# SLIDE 7: Evaluation & Sweet Spot
# ==========================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_background(slide, c_bg)
add_slide_header(slide, "階段五：模型評估與最佳平衡點 (Sweet Spot)")

# Description
desc_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.3), Inches(11.933), Inches(0.8))
tf = desc_box.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "模型越複雜，預測就越好嗎？我們對 10 組特徵選擇演算法的表現進行了特徵個數 k 隨 1 到 5 變化的擬合分析，揭示了精實模型的最佳平衡點 (Sweet Spot)。"
p.font.size = Pt(13)
p.font.color.rgb = c_white
p.font.name = "Microsoft JhengHei"

# Add Left Card
add_card(slide, Inches(0.7), Inches(2.2), Inches(5.8), Inches(4.3), "10 組演算法評估與雙曲線特徵分析", [
    "• 評估對比演算法：",
    "  Pearson 相關、Spearman 相關、F-檢定、互資訊、RFE 遞迴消除、",
    "  SFS 向前步進、SBS 向後步進、Lasso 正則、隨機森林、ElasticNet。",
    "• 模型表現雙趨勢：",
    "  - 測試集決定係數 ($R^2$)：自 k=1 增長，在 k=2 處達到峰值，",
    "    隨後在 k=3..5 處由於引入雜訊特徵而逐漸回落。",
    "  - 測試集預測誤差 (RMSE)：自 k=1 下降，在 k=2 處跌到最低點，",
    "    後續由於維度災難與過擬合，預測誤差再次反彈上升。"
])

# Add Right Card
add_card(slide, Inches(6.8), Inches(2.2), Inches(5.8), Inches(4.3), "Sweet Spot：最佳特徵點 k = 2 (最優解釋力與精實模型)", [
    "• 最佳特徵組合 (Sweet Spot)：",
    "  「研發支出 (R&D Spend)」 與 「市場行銷支出 (Marketing Spend)」。",
    "• 模型指標：",
    "  - 決定係數 $R^2 \\approx 94.74\\%$（相比 5 特徵模型的 93.47% 解釋力更高）。",
    "  - 平均預測誤差 $RMSE \\approx \\$8,198.8$（相比 5 特徵模型的 $9,138.0 誤差更低）。",
    "• 商業啟示：",
    "  多餘的行政支出與地區虛擬特徵在此數據集上實為「噪聲特徵」，",
    "  捨棄它們不仅可以簡化數據採集，更能顯著提升模型的泛化精度！"
], title_color=c_purple)

# Add a small red highlighted text box at the bottom
alert_box = slide.shapes.add_textbox(Inches(0.7), Inches(6.6), Inches(11.933), Inches(0.5))
p = alert_box.text_frame.paragraphs[0]
p.text = "★ 核心結論：精實且高度專注的 2 特徵模型 (R&D + Marketing) 在本商業預測場景中擊敗了全特徵複雜模型。"
p.font.bold = True
p.font.size = Pt(13)
p.font.color.rgb = c_red
p.font.name = "Microsoft JhengHei"

# ==========================================
# SLIDE 8: Deployment
# ==========================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_background(slide, c_bg)
add_slide_header(slide, "階段六：模型部署 (Deployment)")

desc_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.3), Inches(11.933), Inches(1.0))
tf = desc_box.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "模型訓練的最終目的是提供生產環境的決策支援。我們採取 Pipeline 封裝持久化與輕量級 API 部署策略，讓創投經理能在手機或網頁端隨時評估新創標的。"
p.font.size = Pt(14)
p.font.color.rgb = c_white
p.font.name = "Microsoft JhengHei"

# 3 Column cards for deployment
add_card(slide, Inches(0.7), Inches(2.3), Inches(3.7), Inches(4.2), "① 模型持久化封裝\n(Pipeline Serialization)", [
    "• 工具：`joblib` 序列化工具。",
    "• 封裝內容：",
    "  將數據前處理模組 (StandardScaler & OneHotEncoder) 與線性迴歸模型封裝為單一管道物件。",
    "• 導出路徑：",
    "  `startup_profit_model.pkl`",
    "• 價值：",
    "  確保生產環境中輸入數據的標準化計算基準與訓練環境 100% 一致。"
])

add_card(slide, Inches(4.8), Inches(2.3), Inches(3.7), Inches(4.2), "② 輕量級 API 服務\n(FastAPI Integration)", [
    "• API 核心邏輯：",
    "  ```python",
    "  # 啟動時載入 pickle 檔",
    "  model = joblib.load('model.pkl')",
    "  ",
    "  @app.post('/predict')",
    "  def predict(data: StartupInput):",
    "      df = pd.DataFrame([data.dict()])",
    "      pred = model.predict(df)",
    "      return {'profit': pred[0]}",
    "  ```",
    "• 特色：極低延遲，支持高併發查詢。"
])

add_card(slide, Inches(8.9), Inches(2.3), Inches(3.7), Inches(4.2), "③ 商業決策客戶端應用\n(Application Integration)", [
    "• 整合場景：",
    "  可無縫整合至 VC 內部的投資評估儀表板或行動 App。",
    "• 使用場景：",
    "  創投經理在聽取創辦人路演時，輸入研發與行銷預算，即可秒級獲取該公司的合理利潤預估，作為談判與估值參考。"
])

# ==========================================
# SLIDE 9: Business Recommendations
# ==========================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_background(slide, c_bg)
add_slide_header(slide, "基於模型係數的創投 (VC) 投資決策建議")

desc_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.3), Inches(11.933), Inches(1.0))
tf = desc_box.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "模型係數揭示了各項支出與利潤的邊際貢獻（即每多投入 1 美元所能帶來的利潤增減）。以下為基於模型係數的 3 大商業策略建議："
p.font.size = Pt(14)
p.font.color.rgb = c_white
p.font.name = "Microsoft JhengHei"

# 3 Column cards
add_card(slide, Inches(0.7), Inches(2.3), Inches(3.7), Inches(4.2), "① 研發主導戰略\n(R&D-Driven Strategy)", [
    "• 模型係數：研發支出權重為 +0.81。",
    "• 決策意涵：",
    "  每投入 1 美元於研發，利潤將顯著增加 0.81 美元。",
    "• 創投建議：",
    "  新創公司的研發能力與專利深度是獲利的第一推手。VC 應將被投企業的「研發預算佔比」列為最高優先權評估指標，避免投向重行銷輕技術的空殼企業。"
])

add_card(slide, Inches(4.8), Inches(2.3), Inches(3.7), Inches(4.2), "② 精實管理與行政控制\n(Lean Operations)", [
    "• 模型係數：行政支出權重為 -0.03。",
    "• 決策意涵：",
    "  過高的行政與辦公室開支對利潤有負向牽制作用，象徵組織宂餘度大。",
    "• 創投建議：",
    "  督促被投企業落實精實管理。避免新創公司在早期租用豪華辦公室或配置過多非生產性行政人員，減緩資金損耗速度 (Burn Rate)。"
])

add_card(slide, Inches(8.9), Inches(2.3), Inches(3.7), Inches(4.2), "③ 合理市場擴張\n(Rational Marketing)", [
    "• 模型係數：行銷支出權重為 +0.03。",
    "• 決策意涵：",
    "  行銷支出雖對利潤有正向幫助，但邊際效應僅有 0.03（遠低於研發）。",
    "• 創投建議：",
    "  反對早期新創進行大規模補貼或不計後果的燒錢獲客。應在研發產品具備 PMF (產品市場適配) 後，方可有節制地擴大行銷推廣預算。"
])


# Save presentation
prs.save("crisp_dm_regression_slides.pptx")
print("Successfully generated slide deck: crisp_dm_regression_slides.pptx")
